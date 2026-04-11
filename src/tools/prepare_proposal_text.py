# -*- coding: utf-8 -*-
"""
Stage 0: 文本准备 · prepare_proposal_text.py

功能：
- 自动从 src/data/proposals/ 中选择最新的提案文件（或使用 --file 手动指定）
- 接收一个提案文件（目前支持：PDF / DOCX / TXT / MD）
- 尽可能完整地提取每一页/整篇文本
- 对 PDF：
    - 先尝试 pdfplumber 提取文字
    - 如果某页几乎没有文字，自动用 OCR（pytesseract + pdf2image）重试
- 输出：
    - src/data/prepared/<proposal_id>/full_text.txt
    - src/data/prepared/<proposal_id>/pages.json （逐页文本 + 是否用 OCR + 全局偏移）
"""

import os
import json
import argparse
from pathlib import Path

import base64
from io import BytesIO

import pdfplumber
from pdf2image import convert_from_path
from PIL import Image
import pytesseract
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dotenv import load_dotenv
from openai import OpenAI

load_dotenv()


# ========== 配置 ==========

MIN_TEXT_CHARS_PER_PAGE = 30  # 少于这个字数，就认为"文本太少"，尝试 OCR

_PROGRESS_FILE = Path(__file__).resolve().parents[2] / "src" / "data" / "step_progress.json"

def _write_progress(done: int, total: int, pid: str = "") -> None:
    try:
        path = _PROGRESS_FILE.parent / f"step_progress_{pid}.json" if pid else _PROGRESS_FILE
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps({"done": done, "total": total}), encoding="utf-8")
    except Exception:
        pass

# 如果你有中文 OCR 数据，可以加 chi_sim
# 常见配置："eng" / "chi_sim" / "chi_sim+eng"
TESSERACT_LANG = os.getenv("TESS_LANG", "chi_sim+eng")
VISION_MODEL   = os.getenv("VISION_MODEL",  "gpt-4o")
ENABLE_VISION  = os.getenv("ENABLE_VISION", "true").strip().lower() == "true"


# ========== 文件类型识别 ==========

def detect_file_type(path: Path) -> str:
    """
    根据扩展名简单判断文件类型。
    返回: "pdf" / "docx" / "txt"
    """
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix in [".docx", ".doc"]:
        return "docx"
    if suffix in [".txt", ".md"]:
        return "txt"
    if suffix in [".pptx", ".ppt"]:
        return "pptx"
    raise ValueError(f"暂不支持的文件类型: {suffix}")


def find_latest_proposal() -> Path:
    """
    自动在 src/data/proposals/ 下查找"最近修改"的提案文件。
    支持扩展名：.pdf / .docx / .doc / .txt / .md
    """
    base_dir = Path(__file__).resolve().parents[2]
    proposals_dir = base_dir / "src" / "data" / "proposals"

    if not proposals_dir.exists():
        raise FileNotFoundError(f"未找到提案目录: {proposals_dir}")

    candidates = []
    for p in proposals_dir.iterdir():
        if not p.is_file():
            continue
        if p.suffix.lower() not in [".pdf", ".docx", ".doc", ".txt", ".md"]:
            continue
        candidates.append(p)

    if not candidates:
        raise FileNotFoundError(f"提案目录中没有可用文件: {proposals_dir}")

    latest = max(candidates, key=lambda x: x.stat().st_mtime)
    print(f"[INFO] [auto] 选中最新提案文件: {latest}")
    return latest


# ========== 提取函数 ==========

def ocr_page_from_pdf(pdf_path: Path, page_index: int) -> str:
    """
    用 pdf2image + pytesseract 对指定页做 OCR。
    page_index: 0-based
    """
    try:
        images = convert_from_path(
            str(pdf_path),
            first_page=page_index + 1,
            last_page=page_index + 1
        )
    except Exception as e:
        print(f"[WARN] convert_from_path 失败 (page {page_index+1}): {e}")
        return ""

    if not images:
        return ""

    image: Image.Image = images[0]

    try:
        text = pytesseract.image_to_string(image, lang=TESSERACT_LANG)
        return text
    except Exception as e:
        print(f"[WARN] OCR 失败 (page {page_index+1}): {e}")
        return ""


def describe_image_with_vision(pil_image, context_hint: str = "") -> str:
    """Send a PIL Image to the vision LLM; return a plain-text description or '' on any failure."""
    if not ENABLE_VISION:
        return ""
    try:
        buf = BytesIO()
        pil_image.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
        prompt = (
            "You are analyzing a page from a business/research proposal document. "
            "Describe ALL visual content present: charts, graphs (bar/line/pie/scatter), "
            "tables, figures, diagrams, flowcharts. Include key numbers, percentages, "
            "axis labels, legends, and any visible trends or conclusions. "
            "Be concise but complete. Use plain text only. "
            "If there is no meaningful visual content, respond with exactly: "
            "no significant visual content."
        )
        if context_hint:
            prompt += f" Context: {context_hint}"
        _client = OpenAI()
        resp = _client.chat.completions.create(
            model=VISION_MODEL,
            messages=[{"role": "user", "content": [
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {
                    "url": f"data:image/png;base64,{b64}",
                    "detail": "high",
                }},
            ]}],
            max_tokens=512,
        )
        desc = resp.choices[0].message.content.strip()
        return "" if desc.lower().startswith("no significant") else desc
    except Exception as e:
        print(f"[WARN] Vision LLM call failed: {e}", flush=True)
        return ""


def extract_from_pdf(pdf_path: Path, use_ocr: bool = True, pid: str = ""):
    """
    从 PDF 提取逐页文本：
    - 先用 pdfplumber 提取
    - 如果某页文字太少且 use_ocr=True，则使用 OCR

    返回：
      pages_text: list[str]   每页的文本（顺序和页面一致）
      page_sources: list[str] 每页的来源标记："pdf_text" / "ocr" / "empty"
    """
    pages_text = []
    page_sources = []

    pdf_path = Path(pdf_path)

    with pdfplumber.open(pdf_path) as pdf:
        num_pages = len(pdf.pages)
        print(f"[INFO] PDF 页面数: {num_pages}")
        _write_progress(0, num_pages, pid)
        for i, page in enumerate(pdf.pages):
            txt = page.extract_text() or ""
            txt = txt.strip()

            if txt and len(txt) >= MIN_TEXT_CHARS_PER_PAGE:
                pages_text.append(txt)
                page_sources.append("pdf_text")
                print(f"  - 第 {i+1} 页: 使用 pdfplumber 文本，长度 {len(txt)}")
            else:
                if use_ocr:
                    print(f"  - 第 {i+1} 页: 文本太少({len(txt)} chars)，尝试 OCR...")
                    ocr_txt = ocr_page_from_pdf(pdf_path, page_index=i)
                    ocr_txt = (ocr_txt or "").strip()
                    if ocr_txt:
                        pages_text.append(ocr_txt)
                        page_sources.append("ocr")
                        print(f"    -> OCR 成功，长度 {len(ocr_txt)}")
                    else:
                        pages_text.append("")
                        page_sources.append("empty")
                        print("    -> OCR 也没有提取到文本")
                else:
                    pages_text.append(txt)
                    page_sources.append("pdf_text_empty")
                    print(f"  - 第 {i+1} 页: 不启用 OCR，保留空文本")
            _write_progress(i + 1, num_pages, pid)

            # Vision: describe embedded images/charts on this page
            if ENABLE_VISION and page.images:
                print(f"  - 第 {i+1} 页: 发现 {len(page.images)} 个图像，调用视觉模型...", flush=True)
                try:
                    page_imgs = convert_from_path(
                        str(pdf_path), first_page=i + 1, last_page=i + 1, dpi=150
                    )
                    if page_imgs:
                        visual_desc = describe_image_with_vision(
                            page_imgs[0], context_hint=f"Page {i+1} of proposal"
                        )
                        if visual_desc:
                            pages_text[-1] = pages_text[-1] + f"\n[VISUAL CONTENT: {visual_desc}]"
                            print(f"    -> 已附加视觉描述 ({len(visual_desc)} chars)", flush=True)
                except Exception as e:
                    print(f"[WARN] 第 {i+1} 页视觉分析失败: {e}", flush=True)

    return pages_text, page_sources


def extract_from_docx(docx_path: Path):
    """
    从 DOCX 提取文本，并用视觉模型描述内嵌图像。
    这里没有页的概念，就当成一页。
    """
    doc = Document(str(docx_path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    text = "\n".join(paragraphs)

    if ENABLE_VISION and doc.inline_shapes:
        visual_parts = []
        for idx, shape in enumerate(doc.inline_shapes):
            try:
                rId = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                blob = doc.part.related_parts[rId].blob
                pil_image = Image.open(BytesIO(blob))
                desc = describe_image_with_vision(
                    pil_image, context_hint=f"Inline image {idx+1} from DOCX proposal"
                )
                if desc:
                    visual_parts.append(f"[VISUAL CONTENT: {desc}]")
                    print(f"  - DOCX 图像 {idx+1}: 已附加视觉描述", flush=True)
            except Exception as e:
                print(f"[WARN] DOCX 图像 {idx+1} 处理失败: {e}", flush=True)
        if visual_parts:
            text = text + "\n" + "\n".join(visual_parts)

    return [text], ["docx"]


def extract_from_pptx(pptx_path: Path, pid: str = ""):
    """
    从 PPTX 逐张幻灯片提取文本，并对图片形状调用视觉模型描述。
    每张幻灯片当作一页处理。
    """
    prs = Presentation(str(pptx_path))
    pages_text = []
    page_sources = []

    num_slides = len(prs.slides)
    print(f"[INFO] PPTX 幻灯片数: {num_slides}")
    _write_progress(0, num_slides, pid)

    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []
        visual_parts = []

        for shape in slide.shapes:
            # Extract text from any shape with a text frame
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_texts.append(line)

            # Extract embedded pictures for vision analysis
            if ENABLE_VISION:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        blob = shape.image.blob
                        pil_image = Image.open(BytesIO(blob))
                        desc = describe_image_with_vision(
                            pil_image,
                            context_hint=f"Slide {slide_idx + 1} image from PowerPoint proposal"
                        )
                        if desc:
                            visual_parts.append(f"[VISUAL CONTENT: {desc}]")
                            print(f"  - Slide {slide_idx+1} 图像: 已附加视觉描述", flush=True)
                except Exception as e:
                    print(f"[WARN] Slide {slide_idx+1} 图像处理失败: {e}", flush=True)

        slide_text = "\n".join(slide_texts)
        if visual_parts:
            slide_text = slide_text + "\n" + "\n".join(visual_parts)

        pages_text.append(slide_text)
        page_sources.append("pptx")
        print(f"  - Slide {slide_idx+1}: {len(slide_text)} chars", flush=True)
        _write_progress(slide_idx + 1, num_slides, pid)

    return pages_text, page_sources


def extract_from_txt(txt_path: Path):
    """
    从纯文本文件提取。
    同样当成一页。
    """
    content = Path(txt_path).read_text(encoding="utf-8", errors="ignore")
    content = content.strip()
    return [content], ["txt"]


def prepare_text(file_path: Path, proposal_id: str, use_ocr: bool = True):
    """
    核心入口：
    - 识别文件类型
    - 调用对应的提取函数
    - 输出 full_text.txt + pages.json
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(file_path)

    file_type = detect_file_type(file_path)
    print(f"[INFO] 开始提取文本: {file_path} (type={file_type})")

    if file_type == "pdf":
        pages_text, page_sources = extract_from_pdf(file_path, use_ocr=use_ocr, pid=proposal_id)
    elif file_type == "docx":
        pages_text, page_sources = extract_from_docx(file_path)
    elif file_type == "pptx":
        pages_text, page_sources = extract_from_pptx(file_path, pid=proposal_id)
    elif file_type == "txt":
        pages_text, page_sources = extract_from_txt(file_path)
    else:
        raise ValueError(f"未知文件类型: {file_type}")

    full_text = "\n\n".join(pages_text)

    base_dir = Path(__file__).resolve().parents[2]
    out_dir = base_dir / "src" / "data" / "prepared" / proposal_id
    out_dir.mkdir(parents=True, exist_ok=True)

    # 保存 full_text.txt
    full_text_path = out_dir / "full_text.txt"
    full_text_path.write_text(full_text, encoding="utf-8")
    print(f"[OK] full_text.txt 写入完成: {full_text_path} (长度 {len(full_text)} 字符)")

    # 保存 pages.json，增加 global_char_start/global_char_end，方便后续按全局 offset 找页
    pages_data = []
    offset = 0
    for i, (txt, src) in enumerate(zip(pages_text, page_sources)):
        char_len = len(txt)
        page_start = offset
        page_end = offset + char_len
        pages_data.append(
            {
                "page_index": i + 1,
                "source": src,
                "char_len": char_len,
                "global_char_start": page_start,
                "global_char_end": page_end,
                "text": txt,
            }
        )
        # full_text 用 "\n\n" 拼接，所以这里预留 2 个换行符的长度
        offset = page_end + 2

    pages_json_path = out_dir / "pages.json"
    pages_json_path.write_text(json.dumps(pages_data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"[OK] pages.json 写入完成: {pages_json_path}")

    return {
        "proposal_id": proposal_id,
        "file_type": file_type,
        "out_dir": str(out_dir),
        "full_text_path": str(full_text_path),
        "pages_json_path": str(pages_json_path),
        "num_pages": len(pages_text),
    }


# ========== CLI ==========

def main():
    parser = argparse.ArgumentParser(description="Stage 0: 提案文本准备（含 OCR + 自动选最新提案）")
    parser.add_argument("--file", required=False, help="提案文件路径（PDF/DOCX/TXT）。不填则自动选最新提案")
    parser.add_argument("--proposal_id", required=False, help="提案 ID（用于输出目录名，不填则用文件名）")
    parser.add_argument("--no_ocr", action="store_true", help="禁用 OCR（仅调试用）")
    args = parser.parse_args()

    # 1) 确定 file_path
    if args.file:
        file_path = Path(args.file)
        print(f"[INFO] 使用用户指定文件: {file_path}")
    else:
        file_path = find_latest_proposal()

    # 2) 确定 proposal_id
    proposal_id = args.proposal_id or file_path.stem

    # 3) 是否启用 OCR
    use_ocr = not args.no_ocr

    info = prepare_text(file_path, proposal_id, use_ocr=use_ocr)

    print("\n[SUMMARY]")
    print(json.dumps(info, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
